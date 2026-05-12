import argparse
import json
import sys
from pathlib import Path
from pptx import Presentation

def extract_text_from_pptx(pptx_path):
    prs = Presentation(str(pptx_path))
    slides_data = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            # Tables
            table = getattr(shape, "table", None)
            if table is not None:
                cells = []
                for row in table.rows:
                    for cell in row.cells:
                        text = (cell.text or "").strip()
                        if text:
                            cells.append(text)
                if cells:
                    texts.append("\n".join(cells))
                continue
            # Text frames
            if getattr(shape, "has_text_frame", False):
                text = (shape.text or "").strip()
                if text:
                    texts.append(text)
        # Notes (if any)
        notes_text = ""
        try:
            notes_slide = slide.notes_slide
            if notes_slide is not None and getattr(notes_slide, "notes_text_frame", None):
                notes_text = (notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            notes_text = ""
        slides_data.append({
            "slide_number": idx,
            "texts": texts,
            "notes": notes_text,
        })
    return slides_data

def main():
    parser = argparse.ArgumentParser(description="Extract text from a .pptx and output JSON.")
    parser.add_argument("input", type=Path, help="Path to the .pptx file")
    parser.add_argument("-o", "--output", type=Path, help="Output JSON file (defaults to input basename + .json)")
    args = parser.parse_args()

    input_path = args.input
    if not input_path.exists() or not input_path.is_file():
        print("Error: input file does not exist.", file=sys.stderr)
        sys.exit(2)
    # if input_path.suffix.lower() != ".pptx":
    #     print("Error: only .pptx files are supported.", file=sys.stderr)
    #     sys.exit(3)

    data = extract_text_from_pptx(input_path)

    out_path = args.output if args.output else input_path.with_suffix(".json")
    try:
        out_path.parent.mkdir(parents=True, exist_ok=True)
        with out_path.open("w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    except Exception as e:
        print(f"Error writing output: {e}", file=sys.stderr)
        sys.exit(4)

if __name__ == "__main__":
    main()